import argparse
from pptx import Presentation
from pathlib import Path
from rich.console import Console
from rich.progress import track
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide
from jinja2 import Environment
import shutil
import subprocess
from pdf2image import convert_from_path
import re
from rapidfuzz import fuzz
import sys

console = Console()


def main():
    parser = argparse.ArgumentParser(
        description="Convert PPTX files to Typst scripts."
    )

    # Positional argument for input_file
    parser.add_argument(
        "input_file",
        type=Path,
        help="Path to the input PPTX file."
    )

    # Optional argument for output_file
    parser.add_argument(
        "-o", "--output-file",
        type=Path,
        default=None,
        help="Path to the output Typst file."
    )

    # Optional argument for output_dir with default cwd/"assets"
    parser.add_argument(
        "-d", "--output-dir",
        type=Path,
        default=Path.cwd() / "assets",
        help="Directory to save extracted images. Defaults to cwd/'assets'."
    )

    # Optional argument for theme
    parser.add_argument(
        "-t", "--theme",
        type=str,
        default="university",
        help="Theme to use for the Typst slides. Provide path if custom."
    )

    args = parser.parse_args()

    input_file: Path = args.input_file
    output_file: Path = args.output_file
    output_dir: Path = args.output_dir
    theme: str = args.theme

    # Set default output_file if not provided
    if not output_file:
        output_file = input_file.with_suffix(".typ")

    # Verify output_dir is a Path object
    if not isinstance(output_dir, Path):
        console.print("[red]Error:[/red] output_dir is not a Path object.")
        sys.exit(1)

    # Create output_dir if it doesn't exist
    if not output_dir.exists():
        try:
            output_dir.mkdir(parents=True)
            console.print(f"Created output directory: {output_dir}")
        except Exception as e:
            console.print(f"[red]Error creating output directory:[/red] {e}")
            sys.exit(1)

    console.print(f"Reading [bold]{input_file}[/bold]...")
    prs = Presentation(str(input_file))

    # Convert PPTX to images and get image paths
    image_paths = convert_pptx_to_images(
        input_file, output_dir / 'original', output_dir)

    # Process slides
    slides_data = []
    for idx, slide in track(
        enumerate(prs.slides, start=1),
        total=len(prs.slides),
        description="Processing slides",
    ):
        slide_content = process_slide(slide, idx, output_dir)
        # Add the original image path
        if idx <= len(image_paths):
            slide_content['original_image'] = image_paths[idx - 1]
        else:
            slide_content['original_image'] = None
        slides_data.append(slide_content)

    # Generate Typst script
    typst_script = render_typst(slides_data, theme)
    try:
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(typst_script)
        console.print(f"Typst script saved to [bold]{output_file}[/bold]")
    except Exception as e:
        console.print(f"[red]Error writing to output file:[/red] {e}")
        sys.exit(1)


def convert_pptx_to_images(
        input_file: Path, images_dir: Path, output_dir: Path):
    if not shutil.which("libreoffice"):
        console.print(
            "[red]Error:[/red] LibreOffice is not found in PATH.")
        raise sys.exit(1)
    # Ensure images_dir exists
    images_dir.mkdir(parents=True, exist_ok=True)
    # Convert PPTX to PDF
    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        str(input_file),
        "--outdir",
        str(images_dir),
    ]
    console.print("Converting PPTX to PDF using LibreOffice...")
    subprocess.run(command, check=True)
    console.print("Conversion to PDF completed.")
    # Convert PDF to images
    pdf_file = images_dir / input_file.with_suffix(".pdf").name
    if not pdf_file.exists():
        console.print(f"[red]Error:[/red] PDF file not found: {pdf_file}")
        raise sys.exit(1)
    console.print("Converting PDF pages to images...")
    images = convert_from_path(str(pdf_file))
    image_paths = []
    for i, image in enumerate(images):
        image_filename = images_dir / f"slide_{i+1}.png"
        image.save(image_filename, "PNG")
        image_paths.append(str(image_filename.relative_to(
            output_dir.parent)))
    console.print("Conversion to images completed.")
    return image_paths


def escape_special_characters(text):
    # List of special characters to escape
    special_chars = r'([\\{}$#_*])'
    return re.sub(special_chars, r'\\\1', text)


def process_slide(slide: Slide, slide_idx: int, output_dir: Path):
    slide_content = {
        "title": extract_title(slide),
        "content": extract_content_and_links(slide),
        "images": extract_images(slide, slide_idx, output_dir),
        "tables": extract_tables(slide),
    }
    return slide_content


last_title = ""


def extract_title(slide: Slide):
    global last_title
    for shape in slide.shapes:
        if shape.has_text_frame and hasattr(
                shape, 'is_placeholder') and shape.is_placeholder:
            if shape.placeholder_format.type == 1:  # 1 is Title placeholder
                title_text = shape.text.strip()
                # Check for duplicate titles
                if fuzz.ratio(title_text, last_title) < 92:
                    last_title = title_text
                    return title_text
    # If no title placeholder is found, return the first line of text
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    title_text = paragraph.text.strip().split('\n', 1)[0]
                    if fuzz.ratio(title_text, last_title) < 92:
                        last_title = title_text
                        return title_text
    return "Slide"


def extract_content_and_links(slide: Slide):
    body_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and (not hasattr(shape, 'is_placeholder')
                                     or not shape.is_placeholder):
            for paragraph in shape.text_frame.paragraphs:
                para_text = ""
                if paragraph.text.strip():
                    for run in paragraph.runs:
                        run_text = escape_special_characters(run.text)
                        para_text += run_text
                    body_texts.append(para_text)
    return "\n".join(body_texts)


def extract_images(slide: Slide, slide_idx: int, output_dir: Path):
    images = []
    img_count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_count += 1
            image = shape.image
            image_bytes = image.blob
            image_ext = image.ext
            image_filename = output_dir / \
                f"slide{slide_idx+1}_img{img_count}.{image_ext}"
            with open(image_filename, "wb") as img_file:
                img_file.write(image_bytes)
            # Get image dimensions
            width = shape.width.pt  # Points
            height = shape.height.pt
            images.append({
                "path": str(image_filename.relative_to(output_dir.parent)),
                "width": width,
                "height": height
            })
    return images


def extract_tables(slide: Slide):
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table_data = []
            table = shape.table
            for row in table.rows:
                row_data = [
                    escape_special_characters(cell.text.strip())
                    for cell in row.cells
                ]
                table_data.append(row_data)
            if any(len(row) > 1 for row in table_data):
                tables.append(table_data)
    return tables


def render_typst(slides_data, theme):
    # Determine the theme import line
    if Path(theme).exists():
        # It's a file path
        theme_import = f'#import "{theme}": *'
    else:
        # It's a default theme
        theme_import = f'#import themes.{theme}: *'
    front_settings = f'''
#show: {theme}-theme.with(
  aspect-ratio: "16-9",
  footer: self => self.info.institution,
  navigation: "mini-slides",
  config-info(
    title: [Title],
    subtitle: [Subtitle],
    author: [Authors],
    date: datetime.today(),
    institution: [Institution],
  ),
)
#set heading(numbering: numbly("{1}.", default: "1.1"))
#title-slide()

'''
    template_str = '''#import "@preview/touying:0.5.2": *
#import "@preview/cetz:0.2.2"
#import "@preview/numbly:0.1.0": numbly

{{ theme_import }}

{{ front_settings }}

= {{ slides[0].title if slides else 'Presentation' }}

{% for slide in slides %}
== {{ slide.title }}

{{ slide.content }}

{% for image in slide.images %}
#image("{{ image.path }}",width:{{ image.width }}pt,height:{{ image.height }}pt)
{% endfor %}

{% for table in slide.tables %}
#table(
    {% for row in table %}
    [{{ row | join(", ") }}],
    {% endfor %}
)
{% endfor %}

{% if slide.original_image %}
#image("{{ slide.original_image }}")
{% endif %}

{% endfor %}
'''
    env = Environment(
        trim_blocks=True,
        lstrip_blocks=True
    )
    template = env.from_string(template_str)
    typst_script = template.render(slides=slides_data,
                                   theme_import=theme_import,
                                   front_settings=front_settings)
    return typst_script


if __name__ == "__main__":
    main()
